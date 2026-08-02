"""
Extracts raw text from uploaded files: PDF, DOCX, PPTX, TXT, and images (via OCR).
"""
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract


def extract_from_pdf(file_path: str) -> str:
    text_parts = []
    with fitz.open(file_path) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


def extract_from_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def extract_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
    return "\n".join(text_parts)


def extract_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_from_image(file_path: str) -> str:
    image = Image.open(file_path)
    return pytesseract.image_to_string(image)


EXTRACTORS = {
    "pdf": extract_from_pdf,
    "docx": extract_from_docx,
    "pptx": extract_from_pptx,
    "txt": extract_from_txt,
    "image": extract_from_image,
}


def extract_text(file_path: str, file_type: str) -> str:
    extractor = EXTRACTORS.get(file_type)
    if not extractor:
        raise ValueError(f"Unsupported file type for extraction: {file_type}")
    return extractor(file_path)
